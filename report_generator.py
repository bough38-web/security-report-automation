from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from config import PPT_TITLE, PPT_FOOTER, PPT_DATE
import logging

logger = logging.getLogger(__name__)

# --- Design Constants ---
COLOR_PRIMARY = RGBColor(0, 48, 87)       # Navy Blue (KT Brand-like or Corporate)
COLOR_ACCENT = RGBColor(0, 169, 224)      # Cyan/Light Blue
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_GREY = RGBColor(89, 89, 89)
FONT_NAME = "맑은 고딕"  # Korean Font safe bet (Windows/Mac with Office)

def add_title_slide(prs):
    """표지 슬라이드 생성 (커스텀 디자인)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6 = Blank

    # 1. 배경 (전체 Navy)
    bg = slide.shapes.add_shape(
        1, # msoShapeRectangle
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background() # No Line

    # 2. 제목 (White, Bold, Large)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = PPT_TITLE
    p.font.name = FONT_NAME
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # 3. 날짜 (Light Blue)
    date_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"DATE: {PPT_DATE}"
    p.font.name = FONT_NAME
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # 4. 장식용 라인
    line = slide.shapes.add_shape(
        1, Inches(3), Inches(3.8), Inches(4), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

def add_content_slide(prs, title, content):
    """내용 슬라이드 생성 (상단 헤더 + 본문)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank

    # 1. 헤더 바 (Navy)
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()

    # 2. 슬라이드 제목 (White)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"📰 {title} 이슈 요약"
    p.font.name = FONT_NAME
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

    # 3. 본문 박스
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = body_box.text_frame
    tf.word_wrap = True

    # 본문 내용 줄바꿈 처리 및 불릿 포인트 적용
    lines = content.split('\n')
    for line in lines:
        if not line.strip(): continue
        p = tf.add_paragraph()
        p.text = line.strip()
        p.font.name = FONT_NAME
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_GREY
        p.space_after = Pt(10) # 문단 간격
        
        # '-'로 시작하는 목록형 문장은 들여쓰기
        if line.strip().startswith("-") or line.strip().startswith("•"):
             p.level = 0
        else:
             p.level = 0 # 기본

    # 4. 푸터 (Footer)
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.4))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Security Intelligence Center | {PPT_DATE}"
    p.font.name = FONT_NAME
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.RIGHT

def make_ppt(summary_map: dict, path: str):
    """요약본을 기반으로 PowerPoint 보고서 생성 (디자인 적용)"""
    try:
        prs = Presentation()
        
        # 1. 표지
        add_title_slide(prs)

        # 2. 내용 슬라이드
        for keyword, summary in summary_map.items():
            add_content_slide(prs, keyword, summary)

        # 3. 맺음말 (별도 디자인 없이 내용 슬라이드 포맷 활용하되 제목 변경)
        add_content_slide(prs, "보고 대응 및 안내", PPT_FOOTER)

        prs.save(path)
        logger.info(f"PPT 파일 생성 성공: {path}")
    except Exception as e:
        logger.error(f"PPT 생성 중 오류 발생: {e}")
        raise